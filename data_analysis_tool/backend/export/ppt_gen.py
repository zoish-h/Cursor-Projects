"""Generate PowerPoint from analysis."""
import io
from typing import Any, Dict, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt


def _metrics_rows(analysis: Dict[str, Any]) -> List[Tuple[str, str]]:
    dl = analysis.get("dataset_level", {})
    metrics = dl.get("metrics", {}) or {}
    return [(k, str(v)[:150] if not isinstance(v, dict) else str(v)[:150]) for k, v in metrics.items() if k != "explanation_type"]


def build_ppt(analysis: Dict[str, Any]) -> Tuple[bytes, str, str]:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title = slide.shapes.title
    if not title:
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        tf = title.text_frame
        tf.text = "Data Analysis Report"
    else:
        title.text = "Data Analysis Report"
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    sub.text_frame.text = f"{analysis.get('row_count', 0)} rows, {analysis.get('column_count', 0)} columns"

    # Data context
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
    ctx.text_frame.text = "Dataset overview"
    body = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    tf = body.text_frame
    tf.text = f"Columns: {', '.join(analysis.get('columns', [])[:15])}"
    p = tf.add_paragraph()
    p.text = f"Data kind: {analysis.get('dataset_level', {}).get('characterization', {}).get('kind', 'generic')}"
    p2 = tf.add_paragraph()
    p2.text = f"Target: {analysis.get('dataset_level', {}).get('characterization', {}).get('target_column') or 'N/A'}"

    # Metrics
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    mTitle = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    mTitle.text_frame.text = "Key metrics"
    rows = _metrics_rows(analysis)
    top = 1.0
    for name, val in rows[:20]:
        box = slide3.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.4))
        box.text_frame.text = f"{name}: {val}"
        top += 0.35

    # Explainability
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    eTitle = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    eTitle.text_frame.text = "Why these metrics?"
    for i, msg in enumerate(analysis.get("explainability", [])[:5]):
        box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.0 + i * 0.6), Inches(9), Inches(0.6))
        box.text_frame.text = msg

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "analysis_report.pptx"

